import collections
import collections.abc
import os
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract(filename, **kwargs):
    image_count = 0
    presentation = Presentation(filename)
    all_content = {}

    current_slide = 0

    for slide in presentation.slides:
        slide_content = {}

        current_text = 0
        current_image = 0
        current_table = 0

        for shape in slide.shapes:
            # EXTRACT TEXT_PLACEHOLDER
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER or shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                for paragraph in shape.text_frame.paragraphs:
                    text = ''
                    for run in paragraph.runs:
                        text += run.text

                    # CHECK EMPTY TABS
                    if text != '':
                        slide_content.update({f'Text{current_text}': text})
                        current_text += 1

            # EXTRACT IMAGES
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:

                #CREATE IMAGE FOLDER IF NOT EXISTS
                if not os.path.isdir('images'):
                    os.mkdir('images')

                image = shape.image
                image_bytes = image.blob
                image_filename = f'images/image{image_count}.{image.ext}'
                with open(image_filename, 'wb') as f:
                    f.write(image_bytes)
                url = os.path.abspath(image_filename)

                slide_content.update({f'Image{current_image}': url})
                current_image += 1

                image_count += 1

            # EXTRACT TABLES
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table

                table_content = {}
                current_row = 0
                for r in table.rows:
                    ceil_content = []
                    for c in r.cells:
                        ceil_content.append(c.text_frame.text)
                    table_content.update({f'Row{current_row}': ceil_content})
                    current_row += 1
                slide_content.update({f'Table{current_table}': table_content})
                current_table += 1
        all_content.update({f'Slide{current_slide}': slide_content})
        current_slide += 1
    print(all_content)
    with open('data.json', 'w') as outfile:
        json.dump(all_content, outfile, indent=2, ensure_ascii=False)


if __name__ == '__main__':
    extract('test.pptx')

